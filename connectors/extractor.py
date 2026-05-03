"""
AziendaOS - Document Text Extractor

Extracts plain text from binary document files (PDF, DOCX, ODT, XLSX, PPTX, etc.)
for LLM ingestion. Each format is handled by its own specialised library.

Usage:
    from connectors.extractor import extract_text

    with open("contratto.pdf", "rb") as f:
        text = extract_text(f.read(), "contratto.pdf")

    # Returns None on failure (caller decides fallback)
"""

import logging
from io import BytesIO
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)


# ── Lazy-load library imports ───────────────────────────────────────
# Each library is imported only when the corresponding format is used,
# so a missing library doesn't break extraction for other formats.

def _import_pypdf():
    try:
        from pypdf import PdfReader
        return PdfReader
    except ImportError:
        return None


def _import_docx():
    try:
        from docx import Document
        return Document
    except ImportError:
        return None


def _import_pptx():
    try:
        from pptx import Presentation
        return Presentation
    except ImportError:
        return None


def _import_openpyxl():
    try:
        from openpyxl import load_workbook
        return load_workbook
    except ImportError:
        return None


def _import_odf():
    try:
        from odf.opendocument import load as odf_load
        from odf.text import P
        return odf_load, P
    except ImportError:
        return None, None


def _import_mistune():
    try:
        import mistune
        return mistune
    except ImportError:
        return None


# ── Extraction functions ────────────────────────────────────────────

def _extract_pdf(data: bytes) -> Optional[str]:
    PdfReader = _import_pypdf()
    if PdfReader is None:
        return None
    try:
        reader = PdfReader(BytesIO(data))
        pages = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages.append(text.strip())
        return "\n\n".join(pages) if pages else None
    except Exception as exc:
        logger.warning("PDF extraction failed: %s", exc)
        return None


def _extract_docx(data: bytes) -> Optional[str]:
    Document = _import_docx()
    if Document is None:
        return None
    try:
        doc = Document(BytesIO(data))
        parts = []

        # Paragraphs
        for p in doc.paragraphs:
            if p.text.strip():
                parts.append(p.text.strip())

        # Tables
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            parts.append("\n".join(rows))

        return "\n\n".join(parts) if parts else None
    except Exception as exc:
        logger.warning("DOCX extraction failed: %s", exc)
        return None


def _extract_pptx(data: bytes) -> Optional[str]:
    Presentation = _import_pptx()
    if Presentation is None:
        return None
    try:
        prs = Presentation(BytesIO(data))
        parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_parts = [f"--- Slide {slide_num} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_parts.append(shape.text.strip())
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        slide_parts.append(" | ".join(cells))
            parts.append("\n".join(slide_parts))
        return "\n\n".join(parts) if parts else None
    except Exception as exc:
        logger.warning("PPTX extraction failed: %s", exc)
        return None


def _extract_xlsx(data: bytes) -> Optional[str]:
    load_workbook = _import_openpyxl()
    if load_workbook is None:
        return None
    try:
        wb = load_workbook(BytesIO(data), read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_parts = [f"=== Sheet: {sheet_name} ==="]
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                line = "\t".join(cells).strip()
                if line:
                    sheet_parts.append(line)
            parts.append("\n".join(sheet_parts))
        wb.close()
        return "\n\n".join(parts) if parts else None
    except Exception as exc:
        logger.warning("XLSX extraction failed: %s", exc)
        return None


def _extract_odt(data: bytes) -> Optional[str]:
    odf_load, P = _import_odf()
    if odf_load is None:
        return None
    try:
        doc = odf_load(BytesIO(data))
        paras = doc.getElementsByType(P)
        texts = []
        for p in paras:
            # Walk text nodes
            parts = []
            for node in p.childNodes:
                if node.nodeType == node.TEXT_NODE:
                    parts.append(node.data)
                elif hasattr(node, "firstChild") and node.firstChild:
                    if node.firstChild.nodeType == node.TEXT_NODE:
                        parts.append(node.firstChild.data)
            line = "".join(parts).strip()
            if line:
                texts.append(line)
        return "\n\n".join(texts) if texts else None
    except Exception as exc:
        logger.warning("ODT extraction failed: %s", exc)
        return None


def _extract_markdown(data: bytes) -> Optional[str]:
    mistune = _import_mistune()
    text = data.decode("utf-8", errors="replace")
    if mistune:
        try:
            # Strip markdown syntax, keep plain text
            html = mistune.html(text)
            from html.parser import HTMLParser

            class TextStripper(HTMLParser):
                def __init__(self):
                    super().__init__()
                    self.result = []

                def handle_data(self, data):
                    self.result.append(data)

            stripper = TextStripper()
            stripper.feed(html)
            plain = " ".join(stripper.result).strip()
            return plain or text
        except Exception:
            pass
    return text  # fallback: return raw markdown


def _extract_plaintext(data: bytes) -> str:
    """Decode bytes to string. Handles UTF-8 and common fallbacks."""
    for encoding in ("utf-8", "latin-1", "cp1252"):
        try:
            return data.decode(encoding)
        except (UnicodeDecodeError, UnicodeError):
            continue
    return data.decode("utf-8", errors="replace")


# ── Format dispatch table ───────────────────────────────────────────

_FORMAT_HANDLERS = {
    # PDF
    ".pdf": _extract_pdf,
    # Word
    ".doc": _extract_docx,   # legacy .doc may work or fail
    ".docx": _extract_docx,
    # PowerPoint
    ".ppt": _extract_pptx,
    ".pptx": _extract_pptx,
    # Excel
    ".xls": _extract_xlsx,
    ".xlsx": _extract_xlsx,
    # OpenDocument
    ".odt": _extract_odt,
    ".ods": _extract_xlsx,   # ODS -> fallback as XLSX-like (openpyxl may not handle this well)
    ".odp": _extract_pptx,   # ODP -> try PPTX parser
    # Markdown / rich text
    ".md": _extract_markdown,
    ".markdown": _extract_markdown,
    ".rst": _extract_markdown,  # treat as text
    # Plain text
    ".txt": _extract_plaintext,
    ".csv": _extract_plaintext,
    ".log": _extract_plaintext,
    ".json": _extract_plaintext,
    ".xml": _extract_plaintext,
    ".yaml": _extract_plaintext,
    ".yml": _extract_plaintext,
    ".ini": _extract_plaintext,
    ".cfg": _extract_plaintext,
    ".conf": _extract_plaintext,
    ".py": _extract_plaintext,
    ".js": _extract_plaintext,
    ".ts": _extract_plaintext,
    ".html": _extract_plaintext,
    ".css": _extract_plaintext,
    ".sh": _extract_plaintext,
    ".bat": _extract_plaintext,
    ".sql": _extract_plaintext,
    ".env": _extract_plaintext,
    ".gitignore": _extract_plaintext,
    ".dockerfile": _extract_plaintext,
    ".yml": _extract_plaintext,
    ".toml": _extract_plaintext,
}


def extract_text(data: bytes, filename: str) -> Optional[str]:
    """
    Extract plain text from a document's raw bytes.

    Args:
        data: Raw file content as bytes.
        filename: Original filename (used to determine format).

    Returns:
        Extracted plain text, or None if the format is unsupported
        or extraction fails entirely.
    """
    ext = Path(filename).suffix.lower()
    handler = _FORMAT_HANDLERS.get(ext)

    if handler is None:
        logger.debug("No handler for extension '%s' (%s)", ext, filename)
        return None

    try:
        result = handler(data)
        if result and isinstance(result, str) and result.strip():
            return result.strip()
        return None
    except Exception as exc:
        logger.warning("Unexpected error extracting %s: %s", filename, exc)
        return None


def get_supported_extensions() -> list[str]:
    """Return the list of file extensions this extractor can handle."""
    return sorted(_FORMAT_HANDLERS.keys())
